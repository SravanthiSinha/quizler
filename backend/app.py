import json
import random
from flask import Flask, jsonify, request
from flask_cors import CORS
import openai
import os
from werkzeug.utils import secure_filename
import PyPDF2
import docx
from pptx import Presentation

app = Flask(__name__)
CORS(app)
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'docx', 'ppt', 'pptx'}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Set your OpenAI API key
openai.api_key = os.getenv("OPENAI_API_KEY")

@app.route('/')
def hello():
    return jsonify({"message": "Hello, World!"})

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(file_path):
    text = ''
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_ppt(file_path):
    text = ''
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def extract_text_from_file(file_path):
    _, file_extension = os.path.splitext(file_path)
    
    if file_extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension in ['.ppt', '.pptx']:
        return extract_text_from_ppt(file_path)
    elif file_extension == '.docx':
        doc = docx.Document(file_path)
        text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    elif file_extension == '.txt':
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
    else:
        raise ValueError("Unsupported file type")
    
    return text

@app.route('/api/generate_questions', methods=['POST'])
def generate_questions():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        try:
            text = extract_text_from_file(file_path)
            
            # Remove the file after extraction
            os.remove(file_path)
            
            # Get additional parameters from the request
            num_questions = int(request.form.get('num_questions', 5))
            difficulty = request.form.get('difficulty', 'mixed')

            # Introduce randomness by selecting a random portion of the text
            text_length = len(text)
            start_index = random.randint(0, max(0, text_length - 4000))
            selected_text = text[start_index:start_index + 4000]

            prompt = f"""Based on the following text, generate {num_questions} multiple-choice questions with {difficulty} difficulty. 
            Your response must be a valid JSON array of objects, with no additional text before or after. Each object should have the following structure:
            {{
                "qno": <question number>,
                "question": "<the question>",
                "option1": "<first option>",
                "option2": "<second option>",
                "option3": "<third option>",
                "option4": "<fourth option>",
                "answer": "<correct option>",
                "difficulty": "<easy/medium/hard>"
            }}
            Ensure that:
            1. The "answer" field contains the exact text of the correct option.
            2. The "difficulty" field is one of "easy", "medium", or "hard".
            3. The questions are appropriate for the requested difficulty level.
            4. Your entire response is a valid JSON array, starting with '[' and ending with ']'.
            5. Do not include any explanations, comments, or additional text outside the JSON structure.

            Text: {selected_text}
            """

            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant that generates quiz questions based on provided text."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=2000
            )

            content = response.choices[0].message.content.strip()

            if not content:
                return jsonify({"error": "Empty response from API"}), 500

            try:
                # Try to clean the content before parsing
                cleaned_content = content.replace('\n', '').replace('\r', '').strip()
                questions = json.loads(cleaned_content)
                return jsonify(questions)
            except json.JSONDecodeError as json_error:
                return jsonify({
                    "error": "Failed to parse ChatGPT response as JSON",
                    "raw_response": content,
                    "cleaned_response": cleaned_content,
                    "json_error": str(json_error)
                }), 500

        except Exception as e:
            return jsonify({"error": str(e), "type": str(type(e))}), 500
    else:
        return jsonify({"error": "File type not allowed"}), 400

if __name__ == '__main__':
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    app.run(debug=True, host='0.0.0.0', port=8000)
